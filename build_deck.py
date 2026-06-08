#!/usr/bin/env python3
"""Build the team setup how-to deck (python-pptx). One-command 'run apollo setup' workflow."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY="0A2540"; NAVY2="12365E"; CODEBG="0E2238"; GOLD="C9A227"; LIGHT="F4F6F8"
INK="1A2B3C"; MUTE="5B6B7B"; WHITE="FFFFFF"; LINE="D9E0E6"
GREEN="1E7E34"; AMBER="C6890A"; TEAL="17849E"
HFONT="Trebuchet MS"; BFONT="Calibri"; MONO="Courier New"
W, H = 13.333, 7.5

prs=Presentation(); prs.slide_width=Inches(W); prs.slide_height=Inches(H)
BLANK=prs.slide_layouts[6]
def rgb(h): return RGBColor.from_string(h)

def slide(bg=WHITE):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb=rgb(bg); r.line.fill.background(); r.shadow.inherit=False
    return s

def rect(s,x,y,w,h,fill,line=None,lw=1.0,shape=MSO_SHAPE.RECTANGLE,shadow=False):
    sp=s.shapes.add_shape(shape,Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=rgb(fill)
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=rgb(line); sp.line.width=Pt(lw)
    sp.shadow.inherit=False
    if shadow:
        from pptx.oxml.ns import qn
        el=sp._element.spPr
        ef=el.makeelement(qn('a:effectLst'),{}); el.append(ef)
        sh=ef.makeelement(qn('a:outerShdw'),{'blurRad':'70000','dist':'30000','dir':'5400000','rotWithShape':'0'}); ef.append(sh)
        c=sh.makeelement(qn('a:srgbClr'),{'val':'0A2540'}); sh.append(c)
        a=c.makeelement(qn('a:alpha'),{'val':'18000'}); c.append(a)
    return sp

def txt(s,text,x,y,w,h,size=14,color=INK,bold=False,font=BFONT,align=PP_ALIGN.LEFT,
        valign=MSO_ANCHOR.TOP,italic=False,sp_after=4,line_sp=1.0):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=valign
    for m in ('margin_left','margin_right','margin_top','margin_bottom'): setattr(tf,m,0)
    items=text if isinstance(text,list) else [text]
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(sp_after); p.space_before=Pt(0)
        try: p.line_spacing=line_sp
        except: pass
        seg,o=(it if isinstance(it,tuple) else (it,{}))
        r=p.add_run(); r.text=seg; f=r.font
        f.size=Pt(o.get('size',size)); f.name=o.get('font',font)
        f.bold=o.get('bold',bold); f.italic=o.get('italic',italic)
        f.color.rgb=rgb(o.get('color',color))
    return tb

def kicker(s,t,x=0.7,y=0.5,color=GOLD): txt(s,t,x,y,11,0.3,size=12,color=color,bold=True,font=HFONT)
def title(s,t,y=0.82,color=NAVY): txt(s,t,0.7,y,12,0.7,size=30,color=color,bold=True,font=HFONT)
def bullets(s,items,x,y,w,h,size=14,gap=8,color=INK):
    body=[(("•  "+t),{'size':size,'color':color}) for t in items]
    txt(s,body,x,y,w,h,size=size,color=color,font=BFONT,sp_after=gap,line_sp=1.12)
def codecard(s,lines,x,y,w,h,heading=None):
    rect(s,x,y,w,h,CODEBG,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True)
    yy=y+0.18
    if heading: txt(s,heading,x+0.3,yy,w-0.6,0.3,size=11,color=GOLD,bold=True,font=HFONT); yy+=0.42
    txt(s,[(l,{'font':MONO,'size':12.5,'color':"D7E3EF"}) for l in lines],x+0.3,yy,w-0.6,h-(yy-y)-0.15,sp_after=5,line_sp=1.05)

# 1 Title
s=slide(NAVY); rect(s,0,0,0.28,H,GOLD)
kicker(s,"GATHER AI   ·   SALES OPERATIONS",x=0.9,y=1.5)
txt(s,"The Self-Driving",0.9,2.0,11.5,1.0,size=50,color=WHITE,bold=True,font=HFONT,sp_after=0)
txt(s,"Apollo Pipeline",0.9,2.85,11.5,1.0,size=50,color=GOLD,bold=True,font=HFONT,sp_after=0)
txt(s,"Team Setup Guide: stand up your own automated outbound with one command.",0.9,4.05,11,0.5,size=17,color="CADCFC",font=BFONT)
txt(s,"Every morning it finds 20 ICP contacts on your top account, writes a personalized email and call/LinkedIn script for each, and stages them in your Apollo sequence. You review and send.",0.9,4.7,10.8,1.0,size=13.5,color="9FB2CC",font=BFONT,line_sp=1.15)

# 2 What you're standing up
s=slide(LIGHT); kicker(s,"WHAT YOU'RE STANDING UP"); title(s,"One command builds your pipeline, every day")
flow=["Ranked account hopper","Morning driver fires 7am","Research + 20 ICP contacts","Personalized email + scripts","Staged in Apollo + Slack you"]
cw=2.32; gap=0.2; x0=0.7; y=2.5; ch=2.2
for i,f in enumerate(flow):
    x=x0+i*(cw+gap)
    rect(s,x,y,cw,ch,WHITE,line=LINE,lw=1,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True)
    rect(s,x+cw/2-0.28,y+0.3,0.56,0.56,NAVY,shape=MSO_SHAPE.OVAL)
    txt(s,str(i+1),x+cw/2-0.28,y+0.3,0.56,0.56,size=20,color=GOLD,bold=True,font=HFONT,align=PP_ALIGN.CENTER,valign=MSO_ANCHOR.MIDDLE)
    txt(s,f,x+0.18,y+1.05,cw-0.36,ch-1.1,size=14,color=INK,bold=True,font=BFONT,align=PP_ALIGN.CENTER,line_sp=1.05)
    if i<len(flow)-1: txt(s,">",x+cw+gap/2-0.1,y+ch/2-0.3,0.3,0.6,size=24,color=GOLD,bold=True,align=PP_ALIGN.CENTER,valign=MSO_ANCHOR.MIDDLE)
txt(s,"You stay the human gate on send. The system does the finding, writing, and staging.",0.7,5.1,12,0.5,size=14,color=MUTE,italic=True,font=BFONT)

# 3 Prerequisites
s=slide(WHITE); kicker(s,"BEFORE YOU START"); title(s,"What you need connected in Cowork")
reqs=[("Claude","Runs the whole pipeline and your scheduled tasks"),("Apollo","Sequence engine, contacts, sending"),
      ("Clay","Finds and enriches ICP contacts"),("Gmail","Sends + the reply watcher reads it"),("Slack","Where the system reports to you")]
cw=2.32; gap=0.2; x0=0.7; y=2.4; ch=2.4
for i,(t,d) in enumerate(reqs):
    x=x0+i*(cw+gap)
    rect(s,x,y,cw,ch,LIGHT,line=LINE,lw=1,shape=MSO_SHAPE.ROUNDED_RECTANGLE); rect(s,x,y,cw,0.12,GOLD)
    txt(s,t,x+0.2,y+0.4,cw-0.4,0.5,size=18,color=NAVY,bold=True,font=HFONT)
    txt(s,d,x+0.2,y+1.05,cw-0.4,ch-1.2,size=12.5,color=MUTE,font=BFONT,line_sp=1.1)
txt(s,"Missing one? Connect it in Cowork first. Apollo and Clay do the heavy lifting.",0.7,5.3,12,0.5,size=14,color=MUTE,italic=True,font=BFONT)

# 4 How it works
s=slide(LIGHT); kicker(s,"HOW IT WORKS"); title(s,"The architecture, end to end")
boxes=[("HOPPER","Your accounts scored Fit x Signal, ranked best-first"),("MORNING DRIVER","Scheduled task, fires every weekday at 7am"),("APOLLO SEQUENCE","20 contacts staged, paused for your review")]
bw=3.9; gap=0.35; x0=0.7; y=2.2; bh=1.7
for i,(t,d) in enumerate(boxes):
    x=x0+i*(bw+gap)
    rect(s,x,y,bw,bh,NAVY,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True)
    txt(s,t,x+0.25,y+0.28,bw-0.5,0.4,size=15,color=GOLD,bold=True,font=HFONT)
    txt(s,d,x+0.25,y+0.78,bw-0.5,bh-0.9,size=13,color="CADCFC",font=BFONT,line_sp=1.12)
    if i<2: txt(s,">",x+bw+gap/2-0.12,y+bh/2-0.3,0.3,0.6,size=26,color=NAVY,bold=True,align=PP_ALIGN.CENTER,valign=MSO_ANCHOR.MIDDLE)
sub=[("THE GATE decides per account","STRIKE = you approve before send.  NURTURE / QUALIFY = auto."),("RESPONSE WATCHER","Reads Gmail every 15 min, scores replies, Slacks you.")]
bw2=6.0; y2=4.5; bh2=1.5
for i,(t,d) in enumerate(sub):
    x=0.7+i*(bw2+0.33)
    rect(s,x,y2,bw2,bh2,WHITE,line=LINE,lw=1,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True); rect(s,x,y2,0.12,bh2,GOLD)
    txt(s,t,x+0.3,y2+0.25,bw2-0.5,0.4,size=14,color=NAVY,bold=True,font=HFONT)
    txt(s,d,x+0.3,y2+0.75,bw2-0.5,bh2-0.85,size=13,color=MUTE,font=BFONT,line_sp=1.1)

# 5 Hero: one command
s=slide(NAVY); rect(s,0,0,0.28,H,GOLD)
kicker(s,"SETUP IS ONE COMMAND",x=0.9,y=1.15)
txt(s,"Install the skill, then type:",0.9,1.7,11,0.6,size=20,color="CADCFC",font=BFONT)
rect(s,0.9,2.5,11.5,1.4,CODEBG,line=GOLD,lw=1.5,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True)
txt(s,"run apollo setup",0.9,2.5,11.5,1.4,size=44,color=GOLD,bold=True,font=MONO,align=PP_ALIGN.CENTER,valign=MSO_ANCHOR.MIDDLE)
txt(s,"Claude then scaffolds your workspace, finds your Slack and Apollo IDs, writes your config, scores your accounts into the hopper, and registers your morning driver and weekly digest. It stops only for the three things a human must do.",0.9,4.3,11.4,1.4,size=15,color="CADCFC",font=BFONT,line_sp=1.2)
txt(s,[("First time only:  ",{'bold':True,'color':GOLD}),("install the apollo-onboard skill via the Cowork skill installer (it ships in this folder).",{'color':"9FB2CC"})],0.9,5.9,11.4,0.6,size=13,font=BFONT)

# 6 What it does automatically
s=slide(WHITE); kicker(s,'WHAT "RUN APOLLO SETUP" DOES FOR YOU'); title(s,"Automatic, no work from you")
auto=[("Scaffolds the workspace","Creates your state, runs, and hopper files in one go."),
      ("Finds your IDs","Looks up your Slack ID and Apollo mailbox, finds or creates your sequence."),
      ("Writes your config","Fills rep_config and apollo_config with your identity. No copy-paste."),
      ("Scores your accounts","Runs Fit x Signal on your list and builds the ranked hopper."),
      ("Registers your schedule","Sets up the 7am morning driver and the Friday digest."),
      ("Reports back","Tells you exactly what is live and what is left for you.")]
cw=5.85; ch=1.25; gx=0.35; gy=0.25; x0=0.7; y0=2.2
for i,(t,d) in enumerate(auto):
    x=x0+(i%2)*(cw+gx); y=y0+(i//2)*(ch+gy)
    rect(s,x,y,cw,ch,LIGHT,line=LINE,lw=1,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s,x+0.22,y+ch/2-0.16,0.32,0.32,GREEN,shape=MSO_SHAPE.OVAL)
    txt(s,"v",x+0.22,y+ch/2-0.17,0.32,0.34,size=15,color=WHITE,bold=True,font=HFONT,align=PP_ALIGN.CENTER,valign=MSO_ANCHOR.MIDDLE)
    txt(s,t,x+0.7,y+0.16,cw-0.85,0.4,size=14.5,color=NAVY,bold=True,font=HFONT)
    txt(s,d,x+0.7,y+0.58,cw-0.85,ch-0.65,size=12.5,color=MUTE,font=BFONT,line_sp=1.08)

# 7 The 3 things you do
s=slide(LIGHT); kicker(s,"YOUR PART"); title(s,"The only three things you do")
you=[("Paste your accounts","When asked, drop in your assigned account list. Claude scores and ranks them.","30 seconds"),
     ("Create the Apollo fields","One-time, in the Apollo UI: add 4 custom fields and set the day-1 email to use them via the { } picker. Claude points you to the exact clicks.","one-time, ~5 min"),
     ("Click Run now once","In the Scheduled panel, run the driver once and approve the tool prompts. Banks approvals for every future 7am run.","one-time, ~1 min")]
cw=3.86; gap=0.35; x0=0.7; y=2.3; ch=3.4
for i,(t,d,tag) in enumerate(you):
    x=x0+i*(cw+gap)
    rect(s,x,y,cw,ch,WHITE,line=LINE,lw=1,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True)
    rect(s,x+0.3,y+0.3,0.6,0.6,NAVY,shape=MSO_SHAPE.OVAL)
    txt(s,str(i+1),x+0.3,y+0.3,0.6,0.6,size=24,color=GOLD,bold=True,font=HFONT,align=PP_ALIGN.CENTER,valign=MSO_ANCHOR.MIDDLE)
    txt(s,t,x+0.3,y+1.05,cw-0.6,0.7,size=17,color=NAVY,bold=True,font=HFONT,line_sp=1.0)
    txt(s,d,x+0.3,y+1.75,cw-0.6,ch-2.1,size=13,color=INK,font=BFONT,line_sp=1.16)
    txt(s,tag,x+0.3,y+ch-0.5,cw-0.6,0.35,size=12,color=GOLD,bold=True,font=BFONT)
txt(s,"Everything else is automatic. The Apollo fields can't be created by API, so that one is on you, once.",0.7,6.0,12,0.4,size=13,color=MUTE,italic=True,font=BFONT)

# 8 Daily life
s=slide(WHITE); kicker(s,"AFTER SETUP"); title(s,"Your day once it's live")
cards=[("Morning","Claude Slacks you: 20 contacts staged across N accounts, with a one-tap enroll per account."),
       ("You review","Confirm enroll, import that account's apollo-import.csv, then activate the sequence."),
       ("The gate","STRIKE accounts wait for your eyes. NURTURE and QUALIFY are pre-trusted."),
       ("Friday","The weekly digest lands automatically: runs, enrollments, replies, hopper depth.")]
cw=6.0; ch=1.9; gx=0.33; gy=0.3; x0=0.7; y0=2.25
for i,(t,d) in enumerate(cards):
    x=x0+(i%2)*(cw+gx); y=y0+(i//2)*(ch+gy)
    rect(s,x,y,cw,ch,LIGHT,line=LINE,lw=1,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True); rect(s,x,y,0.12,ch,GOLD)
    txt(s,t,x+0.3,y+0.22,cw-0.5,0.4,size=16,color=NAVY,bold=True,font=HFONT)
    txt(s,d,x+0.3,y+0.72,cw-0.55,ch-0.85,size=13.5,color=MUTE,font=BFONT,line_sp=1.15)

# 9 The gate
s=slide(LIGHT); kicker(s,"THE GATE"); title(s,"How each account is handled")
g=[("STRIKE",GREEN,"Fit 70+ and Signal 70+","Prep + stage. You review the content and approve before anything sends. Your top logos."),
   ("NURTURE / QUALIFY",AMBER,"Fit 70+ low signal, or Fit 40-69","Content pre-trusted. One tap to enroll. Keeps volume moving."),
   ("KILL",MUTE,"Either axis under 40","Never worked. Off the queue.")]
cw=3.86; gap=0.35; x0=0.7; y=2.4; ch=3.0
for i,(t,c,sub,d) in enumerate(g):
    x=x0+i*(cw+gap)
    rect(s,x,y,cw,ch,WHITE,line=LINE,lw=1,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True); rect(s,x,y,cw,0.7,c)
    txt(s,t,x+0.25,y,cw-0.5,0.7,size=18,color=WHITE,bold=True,font=HFONT,valign=MSO_ANCHOR.MIDDLE)
    txt(s,sub,x+0.25,y+0.85,cw-0.5,0.6,size=12.5,color=c,bold=True,font=BFONT,line_sp=1.05)
    txt(s,d,x+0.25,y+1.5,cw-0.5,ch-1.6,size=13.5,color=INK,font=BFONT,line_sp=1.18)

# 10 Gotchas
s=slide(WHITE); kicker(s,"GOTCHAS"); title(s,"Four things that will trip you up")
gz=[("Run it once","Bank tool approvals or the automated runs silently do nothing."),
    ("Use the { } picker","Never hand-type merge variables, or emails send literal {{tokens}}."),
    ("Import before you activate","Don't activate a contact before its fields are filled, or it sends blank."),
    ("Enroll + activate need you","Apollo always requires your explicit confirmation to send. By design.")]
cw=6.0; ch=1.9; gx=0.33; gy=0.3; x0=0.7; y0=2.25
for i,(t,d) in enumerate(gz):
    x=x0+(i%2)*(cw+gx); y=y0+(i//2)*(ch+gy)
    rect(s,x,y,cw,ch,LIGHT,line=LINE,lw=1,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True); rect(s,x,y,0.55,ch,NAVY)
    txt(s,"!",x,y,0.55,ch,size=24,color=GOLD,bold=True,font=HFONT,align=PP_ALIGN.CENTER,valign=MSO_ANCHOR.MIDDLE)
    txt(s,t,x+0.75,y+0.25,cw-0.95,0.4,size=15.5,color=NAVY,bold=True,font=HFONT)
    txt(s,d,x+0.75,y+0.72,cw-0.95,ch-0.85,size=13,color=MUTE,font=BFONT,line_sp=1.12)

# 11 Files reference
s=slide(LIGHT); kicker(s,"REFERENCE"); title(s,"What's in the folder")
rows=[("apollo-onboard-skill/","The 'run apollo setup' skill. Install this first."),
      ("TEAM-SETUP.md","Written steps, if you'd rather drive it yourself"),
      ("hopper/setup.py","Scaffolds your workspace"),
      ("hopper/rep_config.json","Your identity: Slack, mailbox, sequence"),
      ("hopper/apollo_config.json","Shared Apollo wiring: sequence, custom fields, import spec"),
      ("hopper/hopper.jsonl","Your ranked account queue"),
      ("hopper/gate.py  ·  digest.py","Queue + gate engine  ·  weekly rollup"),
      ("hopper/MORNING_DRIVER.md","The daily procedure = the scheduled task prompt"),
      ("hopper/APOLLO_SETUP.md","One-time Apollo custom-field wiring")]
y=2.15; rh=0.5
for i,(f,d) in enumerate(rows):
    yy=y+i*rh
    if i%2==0: rect(s,0.7,yy,11.9,rh,WHITE)
    txt(s,f,0.9,yy,5.3,rh,size=13,color=NAVY,bold=True,font=MONO,valign=MSO_ANCHOR.MIDDLE)
    txt(s,d,6.3,yy,6.1,rh,size=13,color=MUTE,font=BFONT,valign=MSO_ANCHOR.MIDDLE)

# 12 Closing
s=slide(NAVY); rect(s,0,0,0.28,H,GOLD)
txt(s,"You're live.",0.9,2.3,11,1.1,size=46,color=WHITE,bold=True,font=HFONT)
txt(s,"Install the skill, type \"run apollo setup\", answer two prompts, do two clicks in Apollo. The 7am driver works your patch from there.",0.9,3.6,10.8,1.1,size=16,color="CADCFC",font=BFONT,line_sp=1.2)
txt(s,"Full written steps live in TEAM-SETUP.md inside the folder.  Questions: ping Peter.",0.9,4.9,10.8,0.6,size=13,color="9FB2CC",italic=True,font=BFONT)

out="/Users/petertosh/Desktop/Claude Cowork/OUTPUTS/Apollo Pipeline Orchestrator/Apollo-Pipeline-Team-Setup-Guide.pptx"
prs.save(out)
print("Saved:",out,"slides:",len(prs.slides._sldIdLst))
